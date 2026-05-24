"""MiGPT 家用智能音箱 — 发布会风格 PPT 生成器 (温暖科技风)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG_DIR = r"C:\Users\17695\Desktop\智能音箱_无实物\presentation\images"
OUTPUT = r"C:\Users\17695\Desktop\智能音箱_无实物\presentation\MiGPT_发布会_v2.pptx"
W, H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide):
    """暖白背景"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xF6, 0xF3, 0xEC)
    shp.line.fill.background()
    # Move to back
    sp = shp._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def _set_alpha(shape, alpha_val):
    """给形状的 solidFill 添加透明度 (alpha) 子元素"""
    import lxml.etree as etree
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    fill_elem = shape.fill._fill._solidFill  # <a:solidFill> (CT_SolidColorFillProperties)
    srgb = fill_elem.find(f'{{{A_NS}}}srgbClr')
    if srgb is not None:
        old_alpha = srgb.find(f'{{{A_NS}}}alpha')
        if old_alpha is not None:
            srgb.remove(old_alpha)
        alpha = etree.SubElement(srgb, f'{{{A_NS}}}alpha')
        alpha.set('val', alpha_val)

def add_img_bg(slide, img_path, darken=True):
    """全屏图片背景 + 可选暗色渐变遮罩"""
    if os.path.exists(img_path):
        shp = slide.shapes.add_picture(img_path, 0, 0, W, H)
        # Move to back
        sp = shp._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
    if darken:
        # Semi-transparent dark overlay
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0x0a, 0x08, 0x06)
        _set_alpha(overlay, '50000')  # 50%
        overlay.line.fill.background()
        # Move overlay above bg
        sp2 = overlay._element
        sp2.getparent().remove(sp2)
        slide.shapes._spTree.insert(3, sp2)

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, italic=False, color=RGBColor(0x2D, 0x20, 0x10), font_name='Inter', align=PP_ALIGN.LEFT, line_space=None):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if line_space:
        p.space_after = Pt(line_space)
    return txBox

def add_label(slide, left, top, text):
    """金色小标签"""
    return add_textbox(slide, left, top, 4, 0.3, text.upper(), font_size=11, bold=True, color=RGBColor(0xe8, 0xa8, 0x4c), font_name='Space Grotesk')

def add_title(slide, left, top, text, size=36):
    """标题"""
    return add_textbox(slide, left, top, 8, 0.8, text, font_size=size, bold=False, color=RGBColor(0x2D, 0x20, 0x10), font_name='DM Serif Display')

def add_card(slide, left, top, width, height, title, body, title_color=RGBColor(0x5C, 0x3D, 0x1E), border_color=None):
    """卡片：白色半透明背景 + 浅边框"""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_alpha(card, '80000')
    card.line.color.rgb = RGBColor(0xE8, 0xE0, 0xD6)
    card.line.width = Pt(0.5)
    # Title
    add_textbox(slide, left + 0.15, top + 0.08, width - 0.3, 0.3, title, font_size=13, bold=True, color=title_color, font_name='Inter')
    # Body
    add_textbox(slide, left + 0.15, top + 0.35, width - 0.3, height - 0.4, body, font_size=10, color=RGBColor(0x8B, 0x73, 0x55), font_name='Inter')

def add_stat_card(slide, left, top, width, height, label, number, suffix='', sub=''):
    """数据统计卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_alpha(card, '80000')
    card.line.color.rgb = RGBColor(0xE8, 0xE0, 0xD6)
    card.line.width = Pt(0.5)
    add_textbox(slide, left + 0.15, top + 0.1, width - 0.3, 0.2, label, font_size=9, color=RGBColor(0x8B, 0x73, 0x55))
    add_textbox(slide, left + 0.15, top + 0.3, width - 0.3, 0.35, f"{number}{suffix}", font_size=24, bold=True, color=RGBColor(0xE8, 0x91, 0x3A), font_name='Space Grotesk')
    if sub:
        add_textbox(slide, left + 0.15, top + 0.65, width - 0.3, 0.2, sub, font_size=8, color=RGBColor(0xC4, 0xB0, 0x9A))

def add_image(slide, path, left, top, width, height):
    """嵌入图片"""
    if os.path.exists(path):
        return slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    return None

print("✅ 正在生成 PPT...")

# ════════════════════ P1 封面 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_img_bg(slide, f"{IMG_DIR}/hero.jpg")
add_label(slide, 0.8, 2.5, "🚀 CRAIC 2026")
txBox = add_textbox(slide, 0.8, 2.9, 10, 1.2, "MiGPT 家用智能音箱", font_size=52, bold=False, color=RGBColor(0xf0, 0xe8, 0xe0), font_name='DM Serif Display')
# Gold gradient on title text - python-pptx doesn't support gradient text easily, so we use gold color
add_textbox(slide, 0.8, 3.9, 10, 0.5, "AI 语音助手  ·  智能家居中枢  ·  桌面信息中心", font_size=18, color=RGBColor(0x90, 0x80, 0x70))
add_textbox(slide, 0.8, 4.6, 10, 0.4, "双引擎驱动  ·  更懂你的家", font_size=14, italic=True, color=RGBColor(0x6b, 0x5c, 0x4c))

# ════════════════════ P2 产品概述 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide2backend.jpg", darken=False)
add_label(slide, 0.6, 0.5, "About")
add_title(slide, 0.6, 0.9, "不只是音箱")
add_textbox(slide, 0.6, 1.6, 10, 0.4, "MiGPT 将传统智能音箱升级为双引擎系统——本地 SOP + 云端大模型，让音箱从被动问答进化为主动式家庭智能管家。", font_size=12, color=RGBColor(0x8B, 0x73, 0x55))
# Left: App UI mockup
add_image(slide, f"{IMG_DIR}/手机端产品图.png", 0.6, 2.2, 3.8, 5.0)
# Right: stat cards
add_stat_card(slide, 7.5, 2.5, 4.8, 1.6, "本地响应速度", "80", "ms", "毫秒级本地 SOP 执行·无需联网")
add_stat_card(slide, 7.5, 4.3, 4.8, 1.6, "日常场景覆盖", "95", "%", "小爱直接执行 + SOP 固化指令")

# ════════════════════ P3 系统架构 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide3backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Architecture")
add_title(slide, 0.6, 0.9, "三道路由架构")
add_textbox(slide, 0.6, 1.5, 10, 0.3, "三层意图分类器判断路径：关键词→语义→LLM 兜底", font_size=11, color=RGBColor(0x8B, 0x73, 0x55))
# Left: image
add_image(slide, f"{IMG_DIR}/slide3backend.png", 0.6, 2.0, 7.0, 4.8)
# Right: three route cards
add_card(slide, 7.8, 2.0, 4.8, 1.3, "🏠 路径A · 小爱优先", "关键词命中（开灯/关窗/天气）→ 直接执行 · 零延迟 · 不计缓存", title_color=RGBColor(0x3a, 0x9e, 0x8c))
add_card(slide, 7.8, 3.5, 4.8, 1.3, "🧠 路径B · 大模型兜底", "非小爱范围 → Qwen2.5 理解 → 回复文本 + 设备指令 JSON · 查缓存", title_color=RGBColor(0xe8, 0xa8, 0x4c))
add_card(slide, 7.8, 5.0, 4.8, 1.3, "🔧 路径C · Reasonix 编程", '触发词"Claude/工作助手" → 安全审批 → 沙箱执行 → 播报', title_color=RGBColor(0x5c, 0xb8, 0x5c))

print("   P1-P3 完成")

# ════════════════════ P4 核心特性 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide4backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Features")
add_title(slide, 0.6, 0.9, "四大核心特性")
# Left: image
add_image(slide, f"{IMG_DIR}/slide4backend.png", 0.6, 1.7, 6.0, 5.0)
add_stat_card(slide, 7.0, 1.7, 5.5, 2.2, "智能路由", "80", "%", "三层意图分类·毫秒级")
add_stat_card(slide, 7.0, 4.2, 5.5, 2.2, "高频缓存", "4", "次·3次学习", "LRU热层 + SQLite冷层")
# Bottom row
add_card(slide, 7.0, 6.4, 5.5, 0.8, "🎯 第三特性 · 主动关怀", "记忆+传感器 → 情境推测 → 主动提醒", title_color=RGBColor(0xE8, 0x91, 0x3A))
# Bottom row of 2x2 - not needed since we have image + 2 cards on right

print("   P4 完成")

# ════════════════════ P5 架构详解 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide5backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Deep Dive")
add_title(slide, 0.6, 0.9, "技术架构详解")
# Left: layer cards
add_card(slide, 0.6, 1.7, 5.5, 1.15, "🌐 前端展示层", "React 18 + TypeScript + Tailwind · 聊天面板 + 设备状态 + 缓存管理", title_color=RGBColor(0x5C, 0x3D, 0x1E))
add_card(slide, 0.6, 3.0, 5.5, 1.15, "⚙️ 后端路由层", "FastAPI + WebSocket · 三层意图分类 · LRU+SQLite 缓存引擎", title_color=RGBColor(0x3a, 0x9e, 0x8c))
add_card(slide, 0.6, 4.3, 5.5, 1.15, "🏗️ 设备模拟层", "VirtualHome 状态字典 · 与米家 API 接口一致 · 场景注入", title_color=RGBColor(0x5c, 0xb8, 0x5c))
add_card(slide, 0.6, 5.6, 5.5, 1.15, "🔒 安全隐私层", "操作审批 · 语音暂停 · 一键清除全部数据", title_color=RGBColor(0xE8, 0x91, 0x3A))
# Right: image
add_image(slide, f"{IMG_DIR}/slide5backend.png", 6.5, 1.7, 6.0, 5.4)

print("   P5 完成")

# ════════════════════ P6 SOP + 缓存 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide6backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Optimization")
add_title(slide, 0.6, 0.9, "越用越聪明")
add_textbox(slide, 0.6, 1.4, 10, 0.3, "自学习 SOP + 高频缓存引擎 · 系统越用越快", font_size=11, color=RGBColor(0x8B, 0x73, 0x55))
# Left: image
add_image(slide, f"{IMG_DIR}/slide6backend.png", 0.6, 2.0, 6.5, 4.8)
# Right: cards
add_card(slide, 7.5, 2.2, 5.0, 2.0, "🔄 自学习 SOP", "观察器记录每次交互→聚类分析≥3次组合→固化为本地快捷指令→秒级执行", title_color=RGBColor(0x5C, 0x3D, 0x1E))
add_card(slide, 7.5, 4.5, 5.0, 2.0, "⚡ 高频缓存引擎", "≥3次计数→自动写入缓存→第4次起秒回·LRU热层+SQLite冷层·7天过期", title_color=RGBColor(0x3a, 0x9e, 0x8c))

print("   P6 完成")

# ════════════════════ P7 主动关怀 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide7backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Proactive Care")
add_title(slide, 0.6, 0.9, "记住你，关心你")
# Left: text content
add_textbox(slide, 0.6, 1.6, 5.5, 0.8, "系统在对话中自动提取你的生活习惯（通勤、起床、偏好），存入知识图谱。情境引擎结合传感器、时间和记忆，在正确时机主动提供帮助。", font_size=11, color=RGBColor(0x8B, 0x73, 0x55))
add_card(slide, 0.6, 2.6, 5.5, 1.5, "🌅 主动提醒示例", '\u201c主人，现在 7:10 了，按您 30 分钟通勤，我帮您拉开窗帘吧？\u201d', title_color=RGBColor(0xE8, 0x91, 0x3A))
add_textbox(slide, 0.6, 4.3, 5.5, 0.6, '触发逻辑：记忆(通勤30分)+时间(7:10)+传感器(卧室有人+窗帘关闭)→主动询问。支持语音指令"关闭智能提醒"随时暂停。', font_size=9, color=RGBColor(0xC4, 0xB0, 0x9A))
# Right: image
add_image(slide, f"{IMG_DIR}/slide7backend.png", 6.8, 1.6, 5.8, 5.0)

print("   P7 完成")

# ════════════════════ P8 技术栈 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide8backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Tech Stack")
add_title(slide, 0.6, 0.9, "技术选型一览")
# Left labels
left_items = [("⚛️ React 18 + TS", "前端框架"), ("🐍 FastAPI", "后端+WebSocket"), ("🦙 Qwen2.5:7B", "本地大模型"), ("🎤 Web Speech API", "语音识别")]
for i, (t, s) in enumerate(left_items):
    y = 2.0 + i * 1.2
    add_card(slide, 0.6, y, 2.8, 0.9, t, s, title_color=RGBColor(0x5C, 0x3D, 0x1E) if i%2==0 else RGBColor(0x3a, 0x9e, 0x8c))
# Center image
add_image(slide, f"{IMG_DIR}/slide8backend.png", 3.7, 1.8, 5.5, 5.0)
# Right labels
right_items = [("🔊 Edge TTS", "语音合成"), ("💾 SQLite + LRU", "双层缓存"), ("🔧 Reasonix", "编程执行"), ("🐳 Docker", "一键部署")]
for i, (t, s) in enumerate(right_items):
    y = 2.0 + i * 1.2
    add_card(slide, 9.5, y, 3.0, 0.9, t, s, title_color=RGBColor(0x5c, 0xb8, 0x5c) if i%2==0 else RGBColor(0xE8, 0x91, 0x3A))

print("   P8 完成")

# ════════════════════ P9 演示场景 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide9backend.png", darken=False)
add_label(slide, 0.6, 0.5, "Demo")
add_title(slide, 0.6, 0.9, "五大演示场景")
# Left: image
add_image(slide, f"{IMG_DIR}/slide9backend.png", 0.6, 1.7, 6.5, 5.0)
# Right: scenario list
scenarios = [
    ("A", "本地秒级响应", '打开卧室灯 -> 0.2s 执行', "#5cb85c"),
    ("B", "复杂语义理解", '我今天好累 -> 大模型推理+设备组合', "#e8a84c"),
    ("C", "Reasonix 编程", '帮我写个脚本 -> 安全审批->执行->播报', "#3a9e8c"),
    ("D", "高频缓存学习", "同一指令 x3 -> 自动缓存 -> 第4次秒回", "#5cb85c"),
    ("E", "主动情境提醒", '记忆+传感器 -> 该出发了要开窗帘吗？', "#fb923c"),
]
for i, (tag, title, desc, color) in enumerate(scenarios):
    y = 1.8 + i * 1.1
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(y), Inches(5.0), Inches(0.9))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_alpha(card, '80000')
    card.line.color.rgb = RGBColor(0xE8, 0xE0, 0xD6)
    card.line.width = Pt(0.5)
    # Colorful left border strip
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(y), Inches(0.06), Inches(0.9))
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(int(color[1:3],16), int(color[3:5],16), int(color[5:7],16))
    border.line.fill.background()
    add_textbox(slide, 7.8, y + 0.08, 4.5, 0.3, f"{tag}  {title}", font_size=13, bold=True, color=RGBColor(0x2D, 0x20, 0x10))
    add_textbox(slide, 7.8, y + 0.45, 4.5, 0.3, desc, font_size=9, color=RGBColor(0x8B, 0x73, 0x55))

print("   P9 完成")

# ════════════════════ P10 尾页 ════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img_bg(slide, f"{IMG_DIR}/slide9backend.png", darken=False)
add_textbox(slide, 0, 3.0, 13.333, 0.04, "━━━━━━━━━━━━━━━━", font_size=14, color=RGBColor(0xE8, 0x91, 0x3A), align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 3.4, 13.333, 1.0, "双引擎驱动 · 更懂你的家", font_size=40, bold=False, color=RGBColor(0x2D, 0x20, 0x10), font_name='DM Serif Display', align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.4, 13.333, 0.4, "MiGPT 家用智能音箱 — CRAIC 2026 参赛项目", font_size=16, color=RGBColor(0x8B, 0x73, 0x55), align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.9, 13.333, 0.3, "纯软件方案 · 三道路由 · 自学习SOP · 渐进式记忆 · 隐私可控", font_size=11, color=RGBColor(0xC4, 0xB0, 0x9A), align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"\n✅ PPT 已保存到: {OUTPUT}")
print(f"   共 {len(prs.slides)} 页")
