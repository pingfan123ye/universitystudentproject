"""第1张幻灯片：封面（v2——不遮挡主体，利用留白区）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


def _add_text_stroke(run, color_hex, alpha_hex):
    """给文字添加描边轮廓，提升白色字在浅色背景上的可读性"""
    rPr = run._r.get_or_add_rPr()
    ln = etree.SubElement(rPr, qn('a:ln'))
    ln.set('w', '635')  # 0.5pt
    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', color_hex)
    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', alpha_hex)


def _for_each_run(slide, fn):
    """对幻灯片所有 textbox 的每个 run 执行 fn"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fn(run)


img_path = r"C:\Users\17695\Desktop\智能音箱_无实物\presentation\images\slide1.png"
output_path = r"C:\Users\17695\Desktop\智能音箱_无实物\presentation\MiGPT_演示.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# === 背景图 ===
slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                          prs.slide_width, prs.slide_height)

# === 主标题（右侧留白区）===
txBox = slide.shapes.add_textbox(Inches(6.2), Inches(0.5), Inches(6.5), Inches(1.3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "MiGPT 家用智能音箱"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)
run.font.name = "Microsoft YaHei"

# === 副标题 ===
txBox2 = slide.shapes.add_textbox(Inches(6.2), Inches(1.8), Inches(6.5), Inches(0.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.LEFT
run2 = p2.add_run()
run2.text = "AI 语音助手  ·  智能家居中枢  ·  桌面信息中心"
run2.font.size = Pt(18)
run2.font.color.rgb = RGBColor(255, 255, 255)
run2.font.name = "Microsoft YaHei"

# === 底部标语 ===
txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(8), Inches(0.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.LEFT
run3 = p3.add_run()
run3.text = "双引擎驱动 · 更懂你的家"
run3.font.size = Pt(16)
run3.font.italic = True
run3.font.color.rgb = RGBColor(255, 255, 255)
run3.font.name = "Microsoft YaHei"

# 给所有文字加黑色半透明描边，确保可读
_for_each_run(slide, lambda r: _add_text_stroke(r, '000000', '30000'))

prs.save(output_path)
print(f"✅ 封面v2已保存: {output_path}")
