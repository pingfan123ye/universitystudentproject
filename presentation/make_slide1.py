"""制作第1张幻灯片：封面"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

img_path = r"C:\Users\17695\Desktop\智能音箱_无实物\presentation\images\slide1.png"
output_path = r"C:\Users\17695\Desktop\智能音箱_无实物\presentation\MiGPT_演示.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# === 背景图（全屏铺满）===
slide.shapes.add_picture(
    img_path,
    Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)

# === 半透明遮罩层（底部，让文字更清晰）===
from pptx.util import Emu
shape = slide.shapes.add_shape(
    1,  # MSO_SHAPE.RECTANGLE
    Inches(0), Inches(4.2),
    prs.slide_width, Inches(3.3)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
shape.fill.fore_color.brightness = 0.6  # 60% 黑色
shape.line.fill.background()

# === 主标题 ===
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.5), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "MiGPT 家用智能音箱"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.LEFT
p.space_after = Pt(8)

# === 副标题 ===
txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(5.7), Inches(10.5), Inches(0.6))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "AI 语音助手 | 智能家居中枢 | 桌面信息中心"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(200, 200, 200)
p2.font.name = "Microsoft YaHei"
p2.alignment = PP_ALIGN.LEFT

# === 底部标语 ===
txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(6.4), Inches(10.5), Inches(0.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "双引擎驱动 · 更懂你的家"
p3.font.size = Pt(16)
p3.font.italic = True
p3.font.color.rgb = RGBColor(160, 160, 160)
p3.font.name = "Microsoft YaHei"
p3.alignment = PP_ALIGN.LEFT

prs.save(output_path)
print(f"✅ 第1张幻灯片已保存到: {output_path}")
